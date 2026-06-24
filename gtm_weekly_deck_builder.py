"""
DSN GTM deck builder.

Two slides, both native editable PowerPoint tables (no images):
  1. Acquisitions and Productivity (by Month)
  2. Deal Creation (by Month)

Shared by the one-time deck and the monthly automation. Pass the data in,
get a .pptx out. Run this file directly to render the current deck.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn

# ---- palette (matches the slide format) -------------------------------
NAVY = RGBColor(0x1D, 0x2A, 0x56)
NAVY_HEADER = RGBColor(0x28, 0x33, 0x5F)
PERI = RGBColor(0x8B, 0x9C, 0xCE)
STAT = RGBColor(0x5D, 0x72, 0xBB)
GRAY = RGBColor(0x7C, 0x82, 0x8F)
ARRBG = RGBColor(0xEE, 0xF1, 0xF9)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
ROWLABEL = RGBColor(0x2B, 0x34, 0x50)
CAPTION = RGBColor(0x9A, 0xA0, 0xAD)
FONT = "Arial"

USABLE_W = 12.13   # inches of table width
LABEL_W = 1.85     # inches for the left label column


def k(v):
    return f"${v / 1000:.1f}K"


def _strip_table_style(table):
    tbl = table._tbl
    tblPr = tbl.find(qn("a:tblPr"))
    if tblPr is not None:
        for el in list(tblPr):
            if el.tag == qn("a:tableStyleId"):
                tblPr.remove(el)
        for a in ("firstRow", "firstCol", "bandRow", "bandCol"):
            tblPr.set(a, "0")


def _no_borders(table):
    for row in table.rows:
        for cell in row.cells:
            tcPr = cell._tc.get_or_add_tcPr()
            for tag in ("a:lnB", "a:lnT", "a:lnR", "a:lnL"):
                ex = tcPr.find(qn(tag))
                if ex is not None:
                    tcPr.remove(ex)
                ln = tcPr.makeelement(qn(tag), {"w": "3175", "cap": "flat"})
                ln.append(ln.makeelement(qn("a:noFill"), {}))
                tcPr.insert(0, ln)


def _cell(cell, text, *, size=10.5, color=NAVY, bold=False, italic=False,
          align=PP_ALIGN.RIGHT, fill=WHITE, wrap=False):
    cell.fill.solid()
    cell.fill.fore_color.rgb = fill
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    cell.margin_left = Inches(0.07)
    cell.margin_right = Inches(0.10)
    cell.margin_top = Inches(0.01)
    cell.margin_bottom = Inches(0.01)
    tf = cell.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    p.text = text
    run = p.runs[0] if p.runs else p.add_run()
    f = run.font
    f.size = Pt(size)
    f.bold = bold
    f.italic = italic
    f.name = FONT
    f.color.rgb = color


def _textbox(slide, left, top, width, height, text, *, size, color,
             bold=False, caps=False):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top),
                                  Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_top = 0
    p = tf.paragraphs[0]
    p.text = text.upper() if caps else text
    r = p.runs[0]
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.name = FONT
    r.font.color.rgb = color
    return tb


def _set_widths(table, n):
    table.columns[0].width = Inches(LABEL_W)
    col_w = Inches((USABLE_W - LABEL_W) / n)
    for i in range(1, n + 1):
        table.columns[i].width = col_w


def _header_row(table, columns):
    _cell(table.cell(0, 0), "", fill=NAVY_HEADER)
    for i, c in enumerate(columns):
        _cell(table.cell(0, i + 1), c, size=11, color=WHITE, bold=True,
              fill=NAVY_HEADER, align=PP_ALIGN.RIGHT)


def _slide_titles(slide, title):
    _textbox(slide, 0.6, 0.42, 12.0, 0.4, "Go-to-Market",
             size=11, color=PERI, bold=True, caps=True)
    _textbox(slide, 0.6, 0.74, 12.1, 0.7, title,
             size=30, color=NAVY, bold=True, caps=True)


def build_acquisitions_slide(prs, columns, segments, period_label, by_word):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _slide_titles(slide, f"Acquisitions and Productivity by {by_word}")

    n = len(columns)
    rows = 1 + len(segments) * 5
    gf = slide.shapes.add_table(rows, n + 1, Inches(0.6), Inches(1.45),
                                Inches(USABLE_W), Inches(5.2))
    table = gf.table
    _strip_table_style(table)
    _set_widths(table, n)

    table.rows[0].height = Inches(0.38)
    _header_row(table, columns)

    r = 1
    for si, seg in enumerate(segments):
        acq, arr, reps = seg["acq"], seg["arr"], seg["reps"]
        tot_acq, tot_arr = sum(acq), sum(arr)
        rep_word = "rep" if reps == 1 else "reps"
        stat = (f"{tot_acq} acquisitions   \u00b7   {k(tot_arr)} ARR   \u00b7   "
                f"{reps} {rep_word} ({period_label})")

        table.rows[r].height = Inches(0.44)
        table.cell(r, 0).merge(table.cell(r, 3))
        table.cell(r, 4).merge(table.cell(r, n))
        _cell(table.cell(r, 0), seg["name"], size=12.5, color=NAVY, bold=True,
              align=PP_ALIGN.LEFT, fill=WHITE, wrap=False)
        _cell(table.cell(r, 4), stat, size=9.5, color=STAT, italic=True,
              align=PP_ALIGN.RIGHT, fill=WHITE)
        r += 1

        table.rows[r].height = Inches(0.30)
        _cell(table.cell(r, 0), "Acquisitions", size=10.5, color=ROWLABEL,
              bold=True, align=PP_ALIGN.LEFT)
        for i, a in enumerate(acq):
            _cell(table.cell(r, i + 1), str(a), size=11, color=NAVY, bold=True)
        r += 1

        table.rows[r].height = Inches(0.28)
        _cell(table.cell(r, 0), "ARR / deal", size=10, color=GRAY,
              align=PP_ALIGN.LEFT)
        for i, a in enumerate(acq):
            _cell(table.cell(r, i + 1), k(arr[i] / a) if a else "\u2014",
                  size=10, color=GRAY)
        r += 1

        table.rows[r].height = Inches(0.28)
        _cell(table.cell(r, 0), "Deals / rep", size=10, color=GRAY,
              align=PP_ALIGN.LEFT)
        for i, a in enumerate(acq):
            _cell(table.cell(r, i + 1), f"{(a/reps if reps else 0):.1f}",
                  size=10, color=GRAY)
        r += 1

        table.rows[r].height = Inches(0.32)
        _cell(table.cell(r, 0), "ARR", size=10.5, color=NAVY, bold=True,
              align=PP_ALIGN.LEFT, fill=ARRBG)
        for i, a in enumerate(acq):
            _cell(table.cell(r, i + 1), k(arr[i]) if a else "\u2014",
                  size=10.5, color=NAVY, bold=True, fill=ARRBG)
        r += 1

    _no_borders(table)
    _textbox(slide, 0.6, 7.12, 12.1, 0.3,
             "Closed-won by close date, calendar month. Current month is "
             "month-to-date. Reconciles to the HubSpot Total Revenue Growth report.",
             size=9, color=CAPTION)
    return slide


def build_creation_slide(prs, columns, rows_data, period_label, by_word):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _slide_titles(slide, f"Deal Creation by {by_word}")

    n = len(columns)
    totals = [sum(rd["counts"][i] for rd in rows_data) for i in range(n)]

    rows = 1 + len(rows_data) + 1
    gf = slide.shapes.add_table(rows, n + 1, Inches(0.6), Inches(1.7),
                                Inches(USABLE_W), Inches(3.6))
    table = gf.table
    _strip_table_style(table)
    _set_widths(table, n)

    table.rows[0].height = Inches(0.46)
    _header_row(table, columns)

    r = 1
    for rd in rows_data:
        table.rows[r].height = Inches(0.62)
        _cell(table.cell(r, 0), rd["name"], size=11.5, color=ROWLABEL,
              bold=True, align=PP_ALIGN.LEFT)
        for i, v in enumerate(rd["counts"]):
            _cell(table.cell(r, i + 1), str(v), size=11.5, color=NAVY)
        r += 1

    table.rows[r].height = Inches(0.62)
    _cell(table.cell(r, 0), "Total", size=11.5, color=NAVY, bold=True,
          align=PP_ALIGN.LEFT, fill=ARRBG)
    for i, v in enumerate(totals):
        _cell(table.cell(r, i + 1), str(v), size=11.5, color=NAVY, bold=True,
              fill=ARRBG)

    _no_borders(table)
    _textbox(slide, 0.6, 4.6, 12.1, 0.5,
             "New Customer and Cloud Conversion / Atlas deals created, by create "
             "date and calendar month. Current month is month-to-date. Excludes "
             "add-on deals and Not Applicable junk.",
             size=9, color=CAPTION)
    return slide


def build_deck(columns, acq_segments, create_rows, out_path,
               period_label="2026 YTD", by_word="Month"):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    build_acquisitions_slide(prs, columns, acq_segments, period_label, by_word)
    build_creation_slide(prs, columns, create_rows, period_label, by_word)
    prs.save(out_path)
    return out_path


if __name__ == "__main__":
    COLUMNS = ["Jan", "Feb", "Mar", "Apr", "May", "Jun"]

    ACQ = [
        {"name": "New Logo Sales (Platform)", "reps": 4,
         "acq": [4, 4, 5, 2, 7, 3],
         "arr": [85764, 80952, 106152, 43478.40, 142951.92, 52884]},
        {"name": "Atlas / Cloud Upgrades", "reps": 4,
         "acq": [2, 3, 2, 3, 4, 4],
         "arr": [53952, 47364, 121728, 93303.96, 95415.84, 203100.84]},
        {"name": "Strategic Accounts (DSOs)", "reps": 1,
         "acq": [2, 1, 3, 0, 2, 3],
         "arr": [30000, 13200, 82800, 0, 22800, 37800]},
    ]

    CREATE = [
        {"name": "New Customer", "counts": [17, 22, 18, 19, 21, 16]},
        {"name": "Cloud Conversion / Atlas", "counts": [9, 11, 18, 21, 13, 12]},
    ]

    build_deck(COLUMNS, ACQ, CREATE, "/home/claude/dsn-gtm-monthly.pptx",
               period_label="2026 YTD", by_word="Month")
    print("built")
